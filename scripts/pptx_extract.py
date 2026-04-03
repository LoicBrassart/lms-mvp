#!/usr/bin/env python3
"""
pptx_extract.py — Extrait le contenu d'un fichier PPTX vers MDX + assets statiques.

Usage:
    python3 scripts/pptx_extract.py <fichier.pptx> [dossier_sortie] [dossier_assets]

Génère :
    <dossier_sortie>/<nom>.mdx       — contenu au format Reveal.js avec frontmatter
    <dossier_assets>/<nom>/images/   — images extraites, nommées slide_NN_img_NN.<ext>
    <dossier_assets>/<nom>/links.csv — liens externes (slide, texte ancre, URL)
"""

import sys
import csv
from datetime import date
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Indices de placeholders considérés comme titres
TITLE_PLACEHOLDER_IDX = {0, 13}


def process_shape(shape, slide_num, img_count, images_dir, links, stem, assets_url_base):
    """
    Traite un shape PPTX et retourne (lignes markdown, nouveau img_count).
    Collecte les liens dans la liste `links` par effet de bord.
    """
    lines = []

    # Images
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        img_count += 1
        image = shape.image
        filename = f"slide_{slide_num:02d}_img_{img_count:02d}.{image.ext}"
        (images_dir / filename).write_bytes(image.blob)
        lines.append(f"![]({assets_url_base}/{stem}/images/{filename})")
        return lines, img_count

    # Groupes : traitement récursif
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            child_lines, img_count = process_shape(
                child, slide_num, img_count, images_dir, links, stem, assets_url_base
            )
            lines.extend(child_lines)
        return lines, img_count

    if not shape.has_text_frame:
        return lines, img_count

    is_title = (
        shape.is_placeholder
        and shape.placeholder_format is not None
        and shape.placeholder_format.idx in TITLE_PLACEHOLDER_IDX
    )

    for para in shape.text_frame.paragraphs:
        parts = []
        for run in para.runs:
            text = run.text
            if not text:
                continue

            if run.hyperlink and run.hyperlink.address:
                url = run.hyperlink.address
                links.append((slide_num, text, url))
                text = f"[{text}]({url})"
            else:
                if run.font.bold and run.font.italic:
                    text = f"***{text}***"
                elif run.font.bold:
                    text = f"**{text}**"
                elif run.font.italic:
                    text = f"*{text}*"

            parts.append(text)

        line = "".join(parts).strip()
        if not line:
            continue

        if is_title:
            line = f"# {line}"
        elif para.level > 0:
            line = "  " * para.level + "- " + line

        lines.append(line)

    return lines, img_count


def extract(pptx_path: str, output_root: str | None = None, assets_root: str | None = None):
    src = Path(pptx_path).resolve()
    if not src.exists():
        sys.exit(f"Erreur : fichier introuvable — {src}")
    if src.suffix.lower() != ".pptx":
        sys.exit(f"Erreur : le fichier doit être un .pptx — {src}")

    stem = src.stem
    base = Path(output_root).resolve() if output_root else src.parent
    mdx_path = base / (stem + ".mdx")

    # Assets sur disque dans public/, URL absolue pour Reveal.js
    assets_fs_root = Path(assets_root).resolve() if assets_root else base / "assets"
    assets_dir = assets_fs_root / stem
    images_dir = assets_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    # Chemin URL absolu : on déduit /assets à partir du nom du dossier racine public/assets
    assets_url_base = "/" + Path(assets_root).as_posix().lstrip("/") if assets_root else "/assets"

    prs = Presentation(src)
    slides_md = []
    all_links = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        img_count = 0

        for shape in slide.shapes:
            lines, img_count = process_shape(
                shape, slide_num, img_count, images_dir, all_links, stem, assets_url_base
            )
            slide_lines.extend(lines)

        slides_md.append("\n".join(slide_lines))

    # Frontmatter
    n_slides = len(prs.slides)
    first_title = next(
        (line[2:].strip() for slide in slides_md for line in slide.splitlines() if line.startswith("# ")),
        stem,
    )
    frontmatter = (
        f'---\n'
        f'title: "{first_title}"\n'
        f'description: "migrated content"\n'
        f'tags: []\n'
        f'publishedAt: {date.today().isoformat()}\n'
        f'level: unknown\n'
        f'duration: {n_slides * 2}\n'
        f'draft: true\n'
        f'---\n\n'
    )

    # <stem>.mdx
    md_content = "\n\n---\n\n".join(slides_md)
    mdx_path.write_text(frontmatter + md_content, encoding="utf-8")

    # links.csv dans le dossier assets
    with open(assets_dir / "links.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["slide", "anchor_text", "url"])
        writer.writerows(all_links)

    n_images = sum(1 for _ in images_dir.iterdir())
    n_links = len(all_links)
    print(f"Extraction terminée →")
    print(f"  {mdx_path}")
    print(f"  {assets_dir}/")
    print(f"  · {n_slides} slides, {n_images} image(s), {n_links} lien(s)")


if __name__ == "__main__":
    if not (2 <= len(sys.argv) <= 4):
        sys.exit("Usage : python3 scripts/pptx_extract.py <fichier.pptx> [dossier_sortie] [dossier_assets]")
    extract(
        sys.argv[1],
        sys.argv[2] if len(sys.argv) >= 3 else None,
        sys.argv[3] if len(sys.argv) == 4 else None,
    )
